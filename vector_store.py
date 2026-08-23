from __future__ import annotations

import os
from io import BytesIO
import docx
import fitz  # PyMuPDF
import openpyxl
import psycopg
import pytesseract
from PIL import Image
from pptx import Presentation
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_postgres import PGVector
from pypdf import PdfReader

load_dotenv()


class ContentExtractionError(Exception):
    """Raised when a file's content can't be extracted as usable text."""


def get_pgvector_connection() -> str:
    direct = os.getenv("PGVECTOR_CONNECTION")
    if direct:
        return direct
    host = os.environ["POSTGRES_HOST"]
    port = os.environ["POSTGRES_PORT"]
    database = os.environ["POSTGRES_DB"]
    user = os.environ["POSTGRES_USER"]
    password = os.environ["POSTGRES_PASSWORD"]
    return f"postgresql+psycopg://{user}:{password}@{host}:{port}/{database}"


def get_collection_name() -> str:
    return os.getenv("PGVECTOR_COLLECTION", "default")


def get_psycopg_connection() -> str:
    return get_pgvector_connection().replace("postgresql+psycopg://", "postgresql://", 1)


def ensure_context_tag_index():
    with psycopg.connect(get_psycopg_connection()) as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                DO $$
                BEGIN
                    IF to_regclass('public.langchain_pg_embedding') IS NOT NULL THEN
                        CREATE INDEX IF NOT EXISTS ix_langchain_pg_embedding_context_tag
                        ON langchain_pg_embedding ((cmetadata->>'context_tag'));
                    END IF;
                END $$;
                """
            )
        conn.commit()


def create_vector_store(embeddings: OpenAIEmbeddings | None = None) -> PGVector:
    vector_embeddings = embeddings or OpenAIEmbeddings()
    store = PGVector(
        embeddings=vector_embeddings,
        collection_name=get_collection_name(),
        connection=get_pgvector_connection(),
        use_jsonb=True,
    )
    ensure_context_tag_index()
    return store


def create_vector_store_from_documents(documents, pre_delete_collection: bool = False) -> PGVector:
    store = PGVector.from_documents(
        documents=documents,
        embedding=OpenAIEmbeddings(),
        collection_name=get_collection_name(),
        connection=get_pgvector_connection(),
        use_jsonb=True,
        pre_delete_collection=pre_delete_collection,
    )
    ensure_context_tag_index()
    return store


def _is_text_unreliable(text: str) -> bool:
    """Detect a PDF text layer that's empty or looks like broken font-encoding
    garbage rather than real content (see docs/ARCHITECTURE.md)."""
    stripped = text.strip()
    if not stripped:
        return True
    recognized = sum(1 for c in stripped if c.isalpha() or c.isspace())
    return (recognized / len(stripped)) < 0.5


def _extract_pdf(raw_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(raw_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    content = "\n".join(pages).strip()
    if not _is_text_unreliable(content):
        return content

    # lang="heb", not "heb+eng": validated against a real Hebrew PDF where
    # the mixed-language model introduced extra misreads on Hebrew glyphs
    # (see docs/ARCHITECTURE.md). PDFs hit this fallback because their
    # existing (single-language, Hebrew) text layer is untrustworthy, so
    # the document's language is already known going in.
    doc = fitz.open(stream=raw_bytes, filetype="pdf")
    ocr_pages = []
    for page in doc:
        pixmap = page.get_pixmap(dpi=300)
        image = Image.open(BytesIO(pixmap.tobytes("png")))
        ocr_pages.append(pytesseract.image_to_string(image, lang="heb").strip())
    doc.close()
    content = "\n".join(ocr_pages).strip()
    if not content:
        raise ContentExtractionError("PDF has no extractable text, even after OCR.")
    return content


def _extract_image(raw_bytes: bytes) -> str:
    # lang="heb+eng": unlike a PDF routed to OCR, an arbitrary dropped-in
    # image (chat/email screenshot, slide export) has no known language
    # ahead of time and may mix Hebrew and English in the same image (see
    # docs/ARCHITECTURE.md) — confirmed against real screenshots of both.
    image = Image.open(BytesIO(raw_bytes))
    return pytesseract.image_to_string(image, lang="heb+eng").strip()


def _extract_docx(raw_bytes: bytes) -> str:
    document = docx.Document(BytesIO(raw_bytes))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    content = "\n".join(paragraphs).strip()
    if not content:
        raise ContentExtractionError("DOCX has no extractable text.")
    return content


def _extract_xlsx(raw_bytes: bytes) -> str:
    workbook = openpyxl.load_workbook(BytesIO(raw_bytes), data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        rows = sheet.iter_rows(values_only=True)
        # Real customer sheets commonly have 1+ title/banner rows (a single
        # non-empty cell) and blank rows before the true header row — take
        # the first row with more than one non-empty cell as the header,
        # rather than assuming row 1 is always it.
        header = None
        for row in rows:
            if sum(1 for cell in row if cell is not None) > 1:
                header = row
                break
        if header is None:
            continue
        header = [str(h) if h is not None else "" for h in header]
        for row in rows:
            pairs = [
                f"{header[i]}: {cell}"
                for i, cell in enumerate(row)
                if cell is not None and i < len(header) and header[i]
            ]
            if pairs:
                lines.append(" | ".join(pairs))
    content = "\n".join(lines).strip()
    if not content:
        raise ContentExtractionError("XLSX has no extractable text.")
    return content


def _extract_pptx(raw_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(raw_bytes))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
    content = "\n".join(lines).strip()
    if not content:
        raise ContentExtractionError("PPTX has no extractable text.")
    return content


def extract_content_from_bytes(raw_bytes: bytes, source: str) -> str:
    extension = source.lower().rsplit(".", 1)[-1] if "." in source else ""

    if extension == "pdf":
        return _extract_pdf(raw_bytes)
    if extension in ("png", "jpg", "jpeg"):
        return _extract_image(raw_bytes)
    if extension == "docx":
        return _extract_docx(raw_bytes)
    if extension == "xlsx":
        return _extract_xlsx(raw_bytes)
    if extension == "pptx":
        return _extract_pptx(raw_bytes)

    try:
        content = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raise ContentExtractionError(
            f"Unsupported file type for '{source}': not UTF-8 text and no parser registered."
        )
    if not content.strip():
        raise ContentExtractionError(f"'{source}' is empty.")
    return content


def _is_hebrew(text: str) -> bool:
    """Heuristic majority-language check, same style as _is_text_unreliable
    above — no LLM call needed just to decide whether translation is
    needed (see docs/ARCHITECTURE.md: translation-at-ingestion)."""
    hebrew_chars = sum(1 for c in text if "א" <= c <= "ת")
    alpha_chars = sum(1 for c in text if c.isalpha())
    return alpha_chars > 0 and (hebrew_chars / alpha_chars) > 0.5


def translate_to_hebrew_if_needed(content: str) -> str:
    """Normalize non-Hebrew document text to Hebrew before chunking/embedding,
    so retrieval isn't split across two embedding neighborhoods by source
    language (see docs/ARCHITECTURE.md: translation-at-ingestion). A no-op
    for documents already majority-Hebrew — most real customer files are, so
    this only spends an LLM call on the minority that need it."""
    if _is_hebrew(content):
        return content
    model = ChatOpenAI(model="gpt-4o")
    response = model.invoke(
        "Translate the following document to Hebrew. Preserve structure, "
        "numbers, names, and formatting as closely as possible. Output only "
        f"the translated text, nothing else.\n\n{content}"
    )
    return response.content
