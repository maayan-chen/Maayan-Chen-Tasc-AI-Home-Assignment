import hashlib
from io import BytesIO
from pathlib import Path

import docx
import fitz  # PyMuPDF
import openpyxl
import pytesseract
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader


class ContentExtractionError(Exception):
    """Raised when a file's content can't be extracted as usable text."""


def _is_text_unreliable(text: str) -> bool:
    """Detect a PDF text layer that's empty or looks like broken font-encoding
    garbage rather than real content (see docs/ARCHITECTURE.md)."""
    stripped = text.strip()
    if not stripped:
        return True
    recognized = sum(1 for c in stripped if c.isalpha() or c.isspace())
    return (recognized / len(stripped)) < 0.5


def _extract_pdf(raw_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(raw_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    content = "\n".join(pages).strip()
    if not _is_text_unreliable(content):
        return content

    # lang="heb", not "heb+eng": validated against a real Hebrew PDF where
    # the mixed-language model introduced extra misreads on Hebrew glyphs
    # (see docs/ARCHITECTURE.md). PDFs hit this fallback because their
    # existing (single-language, Hebrew) text layer is untrustworthy, so
    # the document's language is already known going in.
    doc = fitz.open(stream=raw_bytes, filetype="pdf")
    ocr_pages = []
    for page in doc:
        pixmap = page.get_pixmap(dpi=300)
        image = Image.open(BytesIO(pixmap.tobytes("png")))
        ocr_pages.append(pytesseract.image_to_string(image, lang="heb").strip())
    doc.close()
    content = "\n".join(ocr_pages).strip()
    if not content:
        raise ContentExtractionError("PDF has no extractable text, even after OCR.")
    return content


def _extract_image(raw_bytes: bytes) -> str:
    # lang="heb+eng": unlike a PDF routed to OCR, an arbitrary dropped-in
    # image (chat/email screenshot, slide export) has no known language
    # ahead of time and may mix Hebrew and English in the same image (see
    # docs/ARCHITECTURE.md) — confirmed against real screenshots of both.
    image = Image.open(BytesIO(raw_bytes))
    return pytesseract.image_to_string(image, lang="heb+eng").strip()


def _extract_docx(raw_bytes: bytes) -> str:
    document = docx.Document(BytesIO(raw_bytes))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    content = "\n".join(paragraphs).strip()
    if not content:
        raise ContentExtractionError("DOCX has no extractable text.")
    return content


def _extract_xlsx(raw_bytes: bytes) -> str:
    workbook = openpyxl.load_workbook(BytesIO(raw_bytes), data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        rows = sheet.iter_rows(values_only=True)
        # Real customer sheets commonly have 1+ title/banner rows (a single
        # non-empty cell) and blank rows before the true header row — take
        # the first row with more than one non-empty cell as the header,
        # rather than assuming row 1 is always it.
        header = None
        for row in rows:
            if sum(1 for cell in row if cell is not None) > 1:
                header = row
                break
        if header is None:
            continue
        header = [str(h) if h is not None else "" for h in header]
        for row in rows:
            pairs = [
                f"{header[i]}: {cell}"
                for i, cell in enumerate(row)
                if cell is not None and i < len(header) and header[i]
            ]
            if pairs:
                lines.append(" | ".join(pairs))
    content = "\n".join(lines).strip()
    if not content:
        raise ContentExtractionError("XLSX has no extractable text.")
    return content


def _extract_pptx(raw_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(raw_bytes))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
    content = "\n".join(lines).strip()
    if not content:
        raise ContentExtractionError("PPTX has no extractable text.")
    return content


def extract_content_from_bytes(raw_bytes: bytes, source: str) -> str:
    extension = source.lower().rsplit(".", 1)[-1] if "." in source else ""

    if extension == "pdf":
        return _extract_pdf(raw_bytes)
    if extension in ("png", "jpg", "jpeg"):
        return _extract_image(raw_bytes)
    if extension == "docx":
        return _extract_docx(raw_bytes)
    if extension == "xlsx":
        return _extract_xlsx(raw_bytes)
    if extension == "pptx":
        return _extract_pptx(raw_bytes)

    try:
        content = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raise ContentExtractionError(
            f"Unsupported file type for '{source}': not UTF-8 text and no parser registered."
        )
    if not content.strip():
        raise ContentExtractionError(f"'{source}' is empty.")
    return content


def read_local_files(folder_path: str) -> list[tuple[str, str, str]]:
    folder = Path(folder_path)
    if not folder.exists():
        raise FileNotFoundError(f"Folder not found: {folder_path}")
    if not folder.is_dir():
        raise NotADirectoryError(f"Not a directory: {folder_path}")

    results = []
    for p in folder.rglob("*"):
        if not p.is_file():
            continue
        if p.name.startswith("."):
            continue

        raw_bytes = p.read_bytes()
        try:
            content = extract_content_from_bytes(raw_bytes, source=str(p))
        except ContentExtractionError as e:
            print(f"Skipped {p}: {e}")
            continue

        file_hash = hashlib.sha256(raw_bytes).hexdigest()
        results.append((content, str(p), file_hash))

    return results
